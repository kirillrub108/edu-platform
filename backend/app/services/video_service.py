import hashlib
import os
import re
import shutil
import subprocess
import threading
import zipfile
from collections.abc import Callable
from concurrent.futures import ThreadPoolExecutor, as_completed
from functools import lru_cache
from pathlib import Path

import structlog
from cachetools import TTLCache
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.text.text import TextFrame

from app.config import settings
from app.constants import (
    ENCODE_WORKERS as _ENCODE_WORKERS_DEFAULT,
)
from app.constants import (
    NOWRAP_WIDEN_FACTOR,
    SEGMENT_AUDIO_BITRATE,
    SEGMENT_AUDIO_CHANNELS,
    SEGMENT_FPS,
    SEGMENT_KEYFRAME_SECONDS,
    SLIDE_DPI,
)

# Bundled LibreOffice font-substitution profile; maps Windows/macOS emoji fonts
# (Segoe UI Emoji, Apple Color Emoji) to Noto Color Emoji which is installed in
# the container. Must be copied into the per-run UserInstallation directory so
# LibreOffice picks it up (the -env:UserInstallation flag overrides the default
# ~/.config/libreoffice path, making any static home-dir config irrelevant).
_LO_XCU_SRC = Path(__file__).parent.parent.parent / "lo-emoji-substitution.xcu"

logger = structlog.get_logger()


def _run(cmd: list[str]) -> None:
    result = subprocess.run(cmd, capture_output=True)
    if result.returncode != 0:
        stderr = result.stderr.decode(errors="replace").strip()
        raise RuntimeError(
            f"Command failed (exit {result.returncode}): {' '.join(cmd[:4])}\n{stderr}"
        )


def _slide_number(path: str) -> int:
    m = re.findall(r"(\d+)", Path(path).stem)
    return int(m[-1]) if m else 0


def count_source_slides(full_path: str) -> int | None:
    """Cheap slide/page count for the credit estimate — no rendering.

    PPTX is exact (python-pptx). PDF uses byte heuristics: /Count of the page
    tree, falling back to counting /Type /Page objects; both miss only on
    exotic files with fully compressed object streams. None = can't tell.
    """
    try:
        ext = os.path.splitext(full_path)[1].lower()
        if ext in (".pptx", ".ppt"):
            return len(Presentation(full_path).slides)
        if ext == ".pdf":
            with open(full_path, "rb") as f:
                data = f.read()
            counts = re.findall(rb"/Type\s*/Pages.*?/Count\s+(\d+)", data, re.DOTALL)
            if counts:
                return max(int(c) for c in counts)
            pages = len(re.findall(rb"/Type\s*/Page\b", data)) - len(
                re.findall(rb"/Type\s*/Pages\b", data)
            )
            return pages if pages > 0 else None
    except Exception:
        logger.warning("count_source_slides_failed", path=full_path, exc_info=True)
    return None


def _get_audio_duration(path: str) -> float:
    """Return audio duration in seconds via ffprobe."""
    result = subprocess.run(
        [
            "ffprobe",
            "-v",
            "error",
            "-show_entries",
            "format=duration",
            "-of",
            "default=noprint_wrappers=1:nokey=1",
            path,
        ],
        capture_output=True,
        text=True,
        timeout=30,
    )
    if result.returncode != 0:
        raise RuntimeError(f"ffprobe failed for {path}: {result.stderr.strip()}")
    raw = result.stdout.strip()
    if not raw:
        raise RuntimeError(f"ffprobe returned no duration for {path}")
    return float(raw)


def probe_duration_seconds(path: str) -> float:
    """Public alias of the ffprobe duration probe — works for video too."""
    return _get_audio_duration(path)


def _trim_trailing_silence(src: str, dest: str) -> bool:
    """Write *src* to *dest* with trailing silence removed.

    Threshold: -40 dB, minimum silence run: 0.15 s.
    Returns True when *dest* is valid (≥ 0.1 s) and should be used;
    the caller must fall back to *src* when False is returned.
    """
    r = subprocess.run(
        [
            "ffmpeg",
            "-y",
            "-i",
            src,
            "-af",
            "silenceremove=stop_periods=-1:stop_duration=0.15:stop_threshold=-40dB",
            dest,
        ],
        capture_output=True,
        timeout=60,
    )
    if r.returncode != 0:
        logger.warning("silenceremove_failed", src=src, returncode=r.returncode)
        return False
    try:
        dur = _get_audio_duration(dest)
    except (ValueError, RuntimeError, OSError):
        return False
    if dur < 0.1:
        logger.warning("trimmed_audio_too_short", src=src, duration=round(dur, 3))
        return False
    return True


def _seed_lo_profile(lo_user_dir: str) -> None:
    """Pre-populate LibreOffice UserInstallation with emoji font-substitution config."""
    xcu_dest = Path(lo_user_dir) / "user" / "registrymodifications.xcu"
    xcu_dest.parent.mkdir(parents=True, exist_ok=True)
    if _LO_XCU_SRC.exists():
        shutil.copy2(_LO_XCU_SRC, xcu_dest)
    else:
        logger.warning(
            "lo_xcu_not_found",
            path=str(_LO_XCU_SRC),
        )


# ── PPTX pre-processing: no-wrap text boxes ───────────────────────────────────
# LibreOffice ignores <a:bodyPr wrap="none"> and re-wraps text it had to render
# in a substituted (wider) font, while also declining to grow the spAutoFit box
# the extra line then overflows — so it lands on whatever shape sits below. On
# flush-stacked boxes that reads as lines printed on top of each other. Widening
# the box beforehand removes the reason to wrap; PowerPoint semantics are
# untouched, since wrap="none" means the text never wraps there anyway.
# See docs/DECISIONS.md §55.

_ALGN_FROM_XML: dict[str, PP_ALIGN] = {"ctr": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}


def _text_alignment(text_frame: TextFrame) -> PP_ALIGN:
    """Effective horizontal alignment of a text box.

    OOXML keeps alignment per paragraph, so a box only counts as centred or
    right-aligned when every paragraph stating one agrees; mixed or unset falls
    back to the spec default LEFT, which is also the one direction where
    widening provably cannot move the visible text. When no paragraph says
    anything we still honour the shape-level lstStyle default, otherwise a
    centred box authored that way would drift sideways.
    """
    stated = {p.alignment for p in text_frame.paragraphs if p.alignment is not None}
    if len(stated) == 1:
        only = stated.pop()
        return only if only in (PP_ALIGN.CENTER, PP_ALIGN.RIGHT) else PP_ALIGN.LEFT
    if not stated:
        lvl1 = text_frame._txBody.find(f"{qn('a:lstStyle')}/{qn('a:lvl1pPr')}")
        if lvl1 is not None:
            return _ALGN_FROM_XML.get(lvl1.get("algn", ""), PP_ALIGN.LEFT)
    return PP_ALIGN.LEFT


def _widened_geometry(
    left: int, width: int, slide_cx: int, algn: PP_ALIGN
) -> tuple[int, int] | None:
    """New ``(left, width)`` in EMU for a no-wrap box, or None if it cannot grow.

    The growth direction follows the alignment so the rendered ink stays where
    the author put it: left-aligned grows rightwards, right-aligned leftwards,
    centred evenly around *its own* centre — never the slide's, which would slide
    the text sideways. Every result is clamped inside the slide.
    """
    delta = int(width * (NOWRAP_WIDEN_FACTOR - 1))
    if algn is PP_ALIGN.RIGHT:
        new_left = max(0, left - delta)
        new_width = left + width - new_left
    elif algn is PP_ALIGN.CENTER:
        centre = left + width / 2
        half = min(width / 2 + delta / 2, centre, slide_cx - centre)
        new_left, new_width = int(centre - half), int(2 * half)
    else:
        new_left = left
        new_width = min(width + delta, slide_cx - left)
    if new_width <= width:
        return None
    return new_left, new_width


def _prepare_pptx_for_libreoffice(pptx_path: str, work_dir: str) -> str:
    """Return the file LibreOffice should convert.

    A wrap-relaxed temp copy when the deck has no-wrap boxes worth widening, the
    untouched original otherwise. Never raises: a layout tweak must not be the
    reason a lesson fails to render.
    """
    try:
        prs = Presentation(pptx_path)
        slide_cx = int(prs.slide_width or 0)
    except Exception:
        logger.exception("pptx_prepare_open_failed", path=pptx_path)
        return pptx_path

    if slide_cx <= 0:
        logger.warning("pptx_slide_width_unknown", path=pptx_path)
        return pptx_path

    widened = 0
    for index, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            try:
                if not shape.has_text_frame or shape.text_frame.word_wrap is not False:
                    continue
                if shape.left is None or shape.width is None:
                    continue
                geometry = _widened_geometry(
                    int(shape.left),
                    int(shape.width),
                    slide_cx,
                    _text_alignment(shape.text_frame),
                )
                if geometry is None:
                    # Already flush against the slide edge — nothing to give.
                    logger.warning("nowrap_widen_no_room", slide=index, shape=shape.shape_id)
                    continue
                shape.left, shape.width = geometry
                widened += 1
            except Exception:
                # One malformed shape must not cost the whole presentation.
                logger.exception("nowrap_widen_shape_failed", slide=index)

    if not widened:
        return pptx_path

    dest_dir = os.path.join(work_dir, "_prepared")
    try:
        os.makedirs(dest_dir, exist_ok=True)
        # Same basename — the caller derives the produced PDF's name from the stem.
        dest = os.path.join(dest_dir, os.path.basename(pptx_path))
        prs.save(dest)
    except OSError:
        logger.exception("pptx_prepare_save_failed", path=pptx_path)
        return pptx_path

    logger.info("pptx_nowrap_relaxed", shapes=widened)
    return dest


# ── Font telemetry ────────────────────────────────────────────────────────────
# Pure visibility, not a fix. LibreOffice substitutes missing typefaces silently
# and it is the substitute's metrics that reflow the text, so knowing which
# families real decks actually ask for is what lets the bundled font pack grow
# from evidence instead of guesswork. See docs/DECISIONS.md §56.

_TYPEFACE_RE = re.compile(r'<a:(?:latin|cs|ea) typeface="([^"]*)"')
_FONT_SCAN_PREFIXES = ("ppt/slides/", "ppt/slideLayouts/", "ppt/slideMasters/", "ppt/theme/")


@lru_cache(maxsize=1)
def _installed_font_families() -> frozenset[str]:
    """Lower-cased families fontconfig can resolve. Cached — the container's font
    set cannot change while the process is alive."""
    try:
        result = subprocess.run(
            ["fc-list", ":", "family"],
            capture_output=True,
            text=True,
            errors="replace",
            timeout=30,
        )
    except (OSError, subprocess.SubprocessError):
        logger.warning("fc_list_unavailable")
        return frozenset()
    if result.returncode != 0:
        logger.warning("fc_list_failed", returncode=result.returncode)
        return frozenset()
    families = result.stdout or ""
    return frozenset(
        part.strip().lower()
        for line in families.splitlines()
        for part in line.split(",")
        if part.strip()
    )


def _log_missing_fonts(pptx_path: str) -> None:
    """Warn about typefaces the deck asks for that fontconfig cannot resolve."""
    installed = _installed_font_families()
    if not installed:
        return
    wanted: set[str] = set()
    try:
        with zipfile.ZipFile(pptx_path) as z:
            for name in z.namelist():
                if name.startswith(_FONT_SCAN_PREFIXES) and name.endswith(".xml"):
                    wanted.update(_TYPEFACE_RE.findall(z.read(name).decode("utf8", "replace")))
    except (OSError, zipfile.BadZipFile):
        logger.warning("pptx_font_scan_failed", path=pptx_path)
        return
    # "+mn-lt" & co. are theme references, resolved by the theme's own fontScheme.
    missing = sorted(
        f for f in wanted if f and not f.startswith("+") and f.lower() not in installed
    )
    if missing:
        logger.warning("pptx_fonts_missing", fonts=missing, count=len(missing))


def _pptx_cache_key(pptx_path: str) -> str:
    """Content hash + DPI → stable key for the PNG slide cache."""
    h = hashlib.md5()
    with open(pptx_path, "rb") as f:
        for chunk in iter(lambda: f.read(65536), b""):
            h.update(chunk)
    return f"{h.hexdigest()}_dpi{SLIDE_DPI}"


def _bump_slides_recency(cache_entry_dir: str) -> None:
    """Refresh the cache-entry directory's mtime to "now" on every hit so the
    disk GC (tasks/purge_pipeline) treats it as recently used. Best-effort:
    a read-only mount or a concurrent eviction must never break generation."""
    try:
        os.utime(cache_entry_dir, None)
    except OSError:
        pass


# TTLCache is not thread-safe; the Lock serialises concurrent get/set.
_slides_cache: TTLCache = TTLCache(
    maxsize=settings.SLIDES_CACHE_MAX_SIZE,
    ttl=settings.SLIDES_CACHE_TTL_SECONDS,
)
_slides_cache_lock = threading.Lock()


class VideoService:
    FRAME_RATE = SEGMENT_FPS
    _ENCODE_WORKERS = _ENCODE_WORKERS_DEFAULT

    def extract_slide_texts(self, pptx_path: str) -> list[str]:
        if Path(pptx_path).suffix.lower() not in {".pptx", ".ppt"}:
            logger.info("skip_text_extraction_non_pptx", path=pptx_path)
            return []
        try:
            prs = Presentation(pptx_path)
        except Exception:
            logger.exception("pptx_open_failed", path=pptx_path)
            return []
        slide_texts: list[str] = []
        for slide in prs.slides:
            parts: list[str] = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        parts.append(line)
            slide_texts.append("\n".join(parts))
        return slide_texts

    def convert_pptx_to_images(
        self,
        pptx_path: str,
        output_dir: str,
        cache_dir: str | None = None,
    ) -> list[str]:
        """Convert PPTX/PDF to a sorted list of PNG slide images.

        When *cache_dir* is provided, results are cached by a hash of the
        source file. Repeat conversions of the same PPTX skip LibreOffice and
        pdftoppm entirely and return the cached PNGs (~20-30s saved per hit).
        """
        suffix = Path(pptx_path).suffix.lower()

        # ── cache lookup (PPTX only — PDFs are read directly, no conversion) ──
        if cache_dir and suffix != ".pdf":
            cache_key = _pptx_cache_key(pptx_path)

            with _slides_cache_lock:
                cached_images = _slides_cache.get(cache_key)
            if cached_images:
                logger.info(
                    "slide_cache_hit_memory",
                    cache_key=cache_key[:8],
                    count=len(cached_images),
                )
                # Bump disk recency even on a memory hit: a still-warm entry that
                # never touches disk would look "cold" to the mtime-based cache GC
                # and get evicted, forcing a full LibreOffice re-render next miss.
                _bump_slides_recency(os.path.dirname(cached_images[0]))
                return cached_images

            cached_dir = os.path.join(cache_dir, cache_key)
            cached_images = sorted(
                (str(p) for p in Path(cached_dir).glob("slide-*.png")),
                key=_slide_number,
            )
            if cached_images:
                logger.info(
                    "slide_cache_hit_disk",
                    cache_key=cache_key[:8],
                    count=len(cached_images),
                )
                # mtime bump = "last used" signal the cache GC evicts by (LRU).
                _bump_slides_recency(cached_dir)
                with _slides_cache_lock:
                    _slides_cache[cache_key] = cached_images
                return cached_images

        os.makedirs(output_dir, exist_ok=True)

        if suffix == ".pdf":
            # Use the PDF directly — running it through LibreOffice would corrupt
            # embedded fonts (especially Cyrillic), producing garbled text in the output.
            pdf_path = pptx_path
        else:
            _log_missing_fonts(pptx_path)
            # Convert a wrap-relaxed copy, never the original — the cache key
            # above is computed from the original's bytes and must stay so.
            source_path = _prepare_pptx_for_libreoffice(pptx_path, output_dir)

            pdf_dir = os.path.join(output_dir, "_pdf")
            os.makedirs(pdf_dir, exist_ok=True)

            lo_user_dir = os.path.join(output_dir, "_lo_profile")
            _seed_lo_profile(lo_user_dir)
            _run(
                [
                    "libreoffice",
                    "--headless",
                    f"-env:UserInstallation=file://{lo_user_dir}",
                    "--convert-to",
                    "pdf",
                    "--outdir",
                    pdf_dir,
                    source_path,
                ]
            )

            pdf_name = Path(source_path).stem + ".pdf"
            pdf_path = os.path.join(pdf_dir, pdf_name)
            if not os.path.exists(pdf_path):
                pdfs = list(Path(pdf_dir).glob("*.pdf"))
                if not pdfs:
                    raise RuntimeError(f"LibreOffice produced no PDF from {pptx_path}")
                pdf_path = str(pdfs[0])

        _run(
            [
                "pdftoppm",
                "-png",
                "-r",
                str(SLIDE_DPI),
                "-aa",
                "yes",
                "-aaVector",
                "yes",
                pdf_path,
                os.path.join(output_dir, "slide"),
            ]
        )

        images = sorted(
            (str(p) for p in Path(output_dir).glob("slide-*.png")),
            key=_slide_number,
        )
        if not images:
            raise RuntimeError(f"No slides produced from {pptx_path}")
        logger.info("slides_produced", count=len(images), dpi=SLIDE_DPI)

        # ── populate cache ─────────────────────────────────────────────────────
        if cache_dir and suffix != ".pdf":
            os.makedirs(cached_dir, exist_ok=True)
            for img in images:
                shutil.copy2(img, cached_dir)
            logger.info("slides_cached", count=len(images), cache_key=cache_key[:8])
            with _slides_cache_lock:
                _slides_cache[cache_key] = [
                    os.path.join(cached_dir, Path(img).name) for img in images
                ]

        return images

    def encode_segment(self, idx: int, img: str, aud: str, work_dir: str) -> str:
        """Encode one (image, audio) pair into an MKV segment.

        Thread-safe: all subprocess calls are independent per segment.
        Returns the path to the encoded .mkv file.
        """
        work = Path(work_dir)
        seg_path = str(work / f"_seg_{idx:04d}.mkv")
        trimmed_aud = str(work / f"_aud_{idx:04d}_trim.wav")
        effective_aud = trimmed_aud if _trim_trailing_silence(aud, trimmed_aud) else aud

        duration = _get_audio_duration(effective_aud)
        logger.info("encoding_segment", idx=idx, duration=round(duration, 2))

        _run(
            [
                "ffmpeg",
                "-y",
                "-loop",
                "1",
                "-framerate",
                str(self.FRAME_RATE),
                "-t",
                str(duration),
                "-i",
                img,
                "-i",
                effective_aud,
                # "fast" instead of "medium": ~30% quicker with no visible quality
                # difference for still-image video (tune stillimage suppresses motion
                # estimation anyway, making the preset gap negligible).
                "-c:v",
                "libx264",
                "-tune",
                "stillimage",
                "-preset",
                "fast",
                "-r",
                str(self.FRAME_RATE),
                "-g",
                str(self.FRAME_RATE * SEGMENT_KEYFRAME_SECONDS),
                "-bf",
                "0",
                "-c:a",
                "aac",
                "-ac",
                str(SEGMENT_AUDIO_CHANNELS),
                "-b:a",
                SEGMENT_AUDIO_BITRATE,
                "-ar",
                "48000",
                "-pix_fmt",
                "yuv420p",
                "-vf",
                "scale=trunc(iw/2)*2:trunc(ih/2)*2",
                seg_path,
            ]
        )
        return seg_path

    def concatenate_segments(self, segment_paths: list[str], output_path: str) -> str:
        """Stream-copy pre-encoded MKV segments into the final MP4 without re-encoding."""
        work_dir = Path(output_path).parent
        os.makedirs(work_dir, exist_ok=True)

        list_path = str(work_dir / "_concat_list.txt")
        with open(list_path, "w") as fh:
            for seg in segment_paths:
                fh.write(f"file '{seg}'\n")

        logger.info("concatenating_segments", count=len(segment_paths), output=output_path)
        _run(
            [
                "ffmpeg",
                "-y",
                "-f",
                "concat",
                "-safe",
                "0",
                "-i",
                list_path,
                "-c",
                "copy",
                "-movflags",
                "+faststart",
                output_path,
            ]
        )

        for seg in segment_paths:
            try:
                os.remove(seg)
            except OSError:
                pass
        try:
            os.remove(list_path)
        except OSError:
            pass

        return output_path

    def build_video(
        self,
        image_paths: list[str],
        audio_paths: list[str],
        output_path: str,
        progress_cb: Callable[[int, int], None] | None = None,
    ) -> str:
        """Encode all segments in parallel then concatenate into the final MP4.

        Uses up to _ENCODE_WORKERS concurrent FFmpeg processes. Progress is
        reported via *progress_cb(done, total)* after each segment completes.
        """
        if len(image_paths) != len(audio_paths):
            raise ValueError(
                f"Image/audio count mismatch: {len(image_paths)} vs {len(audio_paths)}"
            )

        os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
        work_dir = str(Path(output_path).parent)
        total = len(image_paths)
        segment_paths: list[str | None] = [None] * total

        with ThreadPoolExecutor(max_workers=self._ENCODE_WORKERS, thread_name_prefix="enc") as pool:
            futures = {
                pool.submit(self.encode_segment, idx, img, aud, work_dir): idx
                for idx, (img, aud) in enumerate(zip(image_paths, audio_paths))
            }
            for future in as_completed(futures):
                idx = futures[future]
                segment_paths[idx] = future.result()
                if progress_cb:
                    progress_cb(sum(1 for p in segment_paths if p), total)

        return self.concatenate_segments(segment_paths, output_path)  # type: ignore[arg-type]


video_service = VideoService()
