"""Shared pytest fixtures for the backend test suite.

The session-scope PG container is brought up once, schema is applied via
`alembic upgrade head`, then every test function runs inside a SAVEPOINT
that is rolled back at teardown. This lets tested code call `await session
.commit()` freely without leaking state between tests.

External services (LLM, TTS, vision, FFmpeg / LibreOffice / pdftoppm,
Redis) are stubbed. Celery is in EAGER mode so `.apply_async` runs
synchronously in-process.
"""

from __future__ import annotations

import io
import os
import subprocess
import uuid
import wave
from collections.abc import AsyncIterator, Iterator
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

import numpy as np
import pytest
import pytest_asyncio


# ── 1. PostgreSQL container ──────────────────────────────────────────────────
# Start the container FIRST and overwrite DATABASE_URL before any `app.*`
# module is imported. `app/tasks/video_pipeline.py` constructs a sync engine
# at import time, so the URL must be valid by then.

from testcontainers.postgres import PostgresContainer


def _pg_container() -> PostgresContainer:
    return PostgresContainer(
        image="postgres:17-alpine",
        username="test_user",
        password="test_password",
        dbname="test_db",
    )


@pytest.fixture(scope="session")
def _postgres() -> Iterator[PostgresContainer]:
    container = _pg_container()
    container.start()
    try:
        yield container
    finally:
        container.stop()


@pytest.fixture(scope="session", autouse=True)
def _set_database_url(_postgres: PostgresContainer, tmp_path_factory: pytest.TempPathFactory) -> Iterator[None]:
    """Patch env BEFORE app modules are imported anywhere."""
    raw_url = _postgres.get_connection_url()  # postgresql+psycopg2://...
    sync_url = raw_url.replace("postgresql+psycopg2", "postgresql+psycopg2")
    async_url = raw_url.replace("postgresql+psycopg2", "postgresql+asyncpg")

    storage = tmp_path_factory.mktemp("storage")
    os.environ["DATABASE_URL"] = async_url
    os.environ["SYNC_DATABASE_URL"] = sync_url
    os.environ["STORAGE_PATH"] = str(storage)
    # Already set via pytest-env, restated for clarity:
    os.environ.setdefault("SECRET_KEY", "test-secret-do-not-use-in-prod")
    os.environ.setdefault("CELERY_TASK_ALWAYS_EAGER", "1")
    os.environ.setdefault("CELERY_TASK_EAGER_PROPAGATES", "1")

    # Force re-read of cached Settings, then rebind the singletons that
    # captured the placeholder URL/path at import time.
    from app.config import get_settings
    get_settings.cache_clear()

    import app.config as _config_mod
    _config_mod.settings = get_settings()

    # 1) async engine in app.database
    from sqlalchemy.ext.asyncio import async_sessionmaker, create_async_engine

    import app.database as _db_mod

    _db_mod.engine = create_async_engine(
        _config_mod.settings.DATABASE_URL,
        echo=False,
        pool_pre_ping=True,
    )
    _db_mod.AsyncSessionLocal = async_sessionmaker(
        _db_mod.engine,
        expire_on_commit=False,
        autoflush=False,
    )

    # 2) sync engine that Celery tasks construct on import
    from sqlalchemy import create_engine
    from sqlalchemy.orm import sessionmaker

    import app.tasks.video_pipeline as _vp_mod

    sync_url = _config_mod.settings.DATABASE_URL.replace("+asyncpg", "+psycopg2")
    _vp_mod.sync_engine = create_engine(sync_url, pool_pre_ping=True)
    _vp_mod.SyncSession = sessionmaker(bind=_vp_mod.sync_engine, expire_on_commit=False)

    # vision_pipeline imports SyncSession by name from video_pipeline — keep
    # its local binding in sync too.
    import app.tasks.vision_pipeline as _vis_mod
    _vis_mod.SyncSession = _vp_mod.SyncSession

    # payment_pipeline also imports SyncSession by name from video_pipeline.
    import app.tasks.payment_pipeline as _pay_mod
    _pay_mod.SyncSession = _vp_mod.SyncSession

    # 2a) usage_service builds its private sync engine lazily on first record;
    # drop anything cached so it re-reads the rebound settings/URL.
    import app.services.usage_service as _usage_mod
    _usage_mod._engine = None
    _usage_mod._SyncSession = None

    # 2b) Celery: even in EAGER mode `task.update_state(...)` writes to the
    # result backend. The default points at Redis (which isn't running in
    # this test env), so swap it for an in-memory backend.
    from app.celery_app import celery_app

    celery_app.conf.update(
        task_always_eager=True,
        task_eager_propagates=True,
        task_store_eager_result=True,
        broker_url="memory://",
        result_backend="cache+memory://",
    )
    # Celery caches the backend on first access — clear so the next
    # `celery_app.backend` resolves the new result_backend.
    celery_app.__dict__.pop("backend", None)

    # 3) storage_service singleton — base_path captured at import
    from app.services.storage_service import StorageService
    import app.services.storage_service as _storage_mod

    _storage_mod.storage_service = StorageService(
        base_path=_config_mod.settings.STORAGE_PATH,
        base_url=_config_mod.settings.BASE_URL,
    )
    # Re-export to other modules that did `from .storage_service import storage_service`
    import app.tasks.video_pipeline as _vp_mod2
    import app.tasks.vision_pipeline as _vis_mod2
    import app.routers.lessons as _lessons_router
    import app.routers.slides as _slides_router
    import app.routers.students as _students_router
    import app.routers.uploads as _uploads_router
    for _mod in (
        _vp_mod2, _vis_mod2, _lessons_router, _slides_router,
        _students_router, _uploads_router,
    ):
        if hasattr(_mod, "storage_service"):
            _mod.storage_service = _storage_mod.storage_service

    yield


# ── 2. Schema bootstrap ──────────────────────────────────────────────────────

@pytest.fixture(scope="session")
def _alembic_upgraded(_set_database_url: None) -> None:
    """Apply migrations once per session. Called by the engine fixture."""
    from alembic import command
    from alembic.config import Config

    # Resolve alembic.ini relative to this file: backend/tests/conftest.py
    # → backend/alembic.ini
    backend_dir = Path(__file__).resolve().parent.parent
    cfg = Config(str(backend_dir / "alembic.ini"))
    cfg.set_main_option("script_location", str(backend_dir / "alembic"))
    cfg.set_main_option(
        "sqlalchemy.url",
        os.environ["DATABASE_URL"].replace("+asyncpg", "+psycopg2"),
    )
    command.upgrade(cfg, "head")


# ── 3. Async engine + per-test session with SAVEPOINT rollback ──────────────

@pytest_asyncio.fixture(scope="session")
async def _async_engine(_alembic_upgraded: None) -> AsyncIterator[Any]:
    from sqlalchemy.ext.asyncio import create_async_engine

    engine = create_async_engine(
        os.environ["DATABASE_URL"],
        echo=False,
        pool_pre_ping=True,
    )
    try:
        yield engine
    finally:
        await engine.dispose()


@pytest_asyncio.fixture()
async def db_session(_async_engine: Any) -> AsyncIterator[Any]:
    """Function-scope AsyncSession joined to an external transaction.

    Canonical SQLAlchemy 2.0 async recipe (see "Joining a Session into an
    External Transaction" in the docs): outer txn + SAVEPOINT, the listener
    on `sync_session` reopens a fresh SAVEPOINT on the underlying
    sync_connection each time an inner txn ends. Anything the tested code
    commits is contained in the SAVEPOINT and discarded when the outer
    transaction is rolled back at teardown.
    """
    from sqlalchemy import event
    from sqlalchemy.ext.asyncio import AsyncSession

    async with _async_engine.connect() as connection:
        await connection.begin()
        await connection.begin_nested()

        session = AsyncSession(
            bind=connection, expire_on_commit=False, autoflush=False
        )

        @event.listens_for(session.sync_session, "after_transaction_end")
        def _end_savepoint(sess, transaction):  # type: ignore[no-redef]
            if connection.closed:
                return
            if not connection.sync_connection.in_nested_transaction():
                connection.sync_connection.begin_nested()

        try:
            yield session
        finally:
            event.remove(
                session.sync_session, "after_transaction_end", _end_savepoint
            )
            await session.close()
            await connection.rollback()


# ── 4. FastAPI app + httpx AsyncClient with dependency overrides ────────────

@pytest_asyncio.fixture()
async def app(db_session: Any) -> AsyncIterator[Any]:
    """FastAPI app instance with get_db and get_redis overridden, rate
    limit disabled, and lifespan bypassed (we already applied migrations).
    """
    import fakeredis.aioredis

    from app.database import get_db
    from app.main import app as fastapi_app
    from app.redis_client import get_redis

    # Disable slowapi for the duration of this fixture so per-route
    # decorators don't 429 our parametrized tests.
    fastapi_app.state.limiter.enabled = False

    fake = fakeredis.aioredis.FakeRedis(decode_responses=True)

    async def _override_get_db() -> AsyncIterator[Any]:
        # Each request inside one test reuses the same SAVEPOINT-bound
        # session, so route-level commits still rollback cleanly.
        yield db_session

    async def _override_get_redis() -> Any:
        return fake

    fastapi_app.dependency_overrides[get_db] = _override_get_db
    fastapi_app.dependency_overrides[get_redis] = _override_get_redis

    try:
        yield fastapi_app
    finally:
        fastapi_app.dependency_overrides.clear()
        await fake.aclose()
        fastapi_app.state.limiter.enabled = True


@pytest_asyncio.fixture()
async def client(app: Any) -> AsyncIterator[Any]:
    from httpx import ASGITransport, AsyncClient

    async def _auto_csrf(request: Any) -> None:
        """For state-changing methods, read csrf_token from the outgoing Cookie
        header and inject it as X-CSRF-Token so tests don't need to do this
        manually on every POST/PUT/PATCH/DELETE call."""
        if request.method not in ("POST", "PUT", "PATCH", "DELETE"):
            return
        cookie_str = request.headers.get("cookie", "")
        for part in cookie_str.split(";"):
            key, _, val = part.strip().partition("=")
            if key == "csrf_token" and val:
                # httpx Headers._list stores 3-tuples (raw_key, normalized_key, value).
                # Append in-place so the same Headers object is mutated — replacing
                # the attribute would target the wrong reference.
                request.headers._list.append(
                    (b"x-csrf-token", b"x-csrf-token", val.encode("latin-1"))
                )
                break

    transport = ASGITransport(app=app)
    async with AsyncClient(
        transport=transport,
        base_url="http://testserver",
        event_hooks={"request": [_auto_csrf]},
    ) as ac:
        yield ac


# ── 5. User + JWT helpers ────────────────────────────────────────────────────

@pytest_asyncio.fixture()
async def teacher_user(db_session: Any) -> Any:
    from app.models.user import User, UserRole
    from app.services.auth_service import hash_password

    user = User(
        email=f"teacher-{uuid.uuid4().hex[:8]}@example.com",
        hashed_password=hash_password("teacher-pass-123"),
        full_name="Teacher One",
        role=UserRole.teacher,
        is_active=True,
        email_verified=True,
    )
    db_session.add(user)
    await db_session.commit()
    await db_session.refresh(user)
    return user


@pytest_asyncio.fixture()
async def student_user(db_session: Any) -> Any:
    from app.models.user import User, UserRole
    from app.services.auth_service import hash_password

    user = User(
        email=f"student-{uuid.uuid4().hex[:8]}@example.com",
        hashed_password=hash_password("student-pass-123"),
        full_name="Student One",
        role=UserRole.student,
        is_active=True,
    )
    db_session.add(user)
    await db_session.commit()
    await db_session.refresh(user)
    return user


_TEST_CSRF = "test-csrf-fixed-value"


def _bearer(user: Any) -> dict[str, str]:
    """Return a cookie dict with access_token + csrf_token for test requests.
    The client fixture's event hook auto-adds X-CSRF-Token for POST/PUT/PATCH/DELETE."""
    from app.services.auth_service import create_access_token

    token, _jti, _exp = create_access_token(user)
    return {"access_token": token, "csrf_token": _TEST_CSRF}


@pytest.fixture()
def teacher_token(teacher_user: Any) -> dict[str, str]:
    return _bearer(teacher_user)


@pytest.fixture()
def student_token(student_user: Any) -> dict[str, str]:
    return _bearer(student_user)


# ── 6. External-service mocks ────────────────────────────────────────────────


@pytest.fixture(autouse=True)
def mock_send_email(monkeypatch: pytest.MonkeyPatch) -> Any:
    """Stub the send_email task's enqueue so no real email provider is hit.

    Celery runs in EAGER mode, so an un-stubbed `send_email.delay(...)` would
    execute the task body in-process and call Resend. Patching `.delay` on the
    shared task object covers every caller (auth router, video pipeline), since
    they all reference the same registered task. Returns the Mock so tests can
    assert enqueue calls / read kwargs.
    """
    from unittest.mock import MagicMock

    from app.tasks import email_pipeline

    m = MagicMock(name="send_email.delay")
    monkeypatch.setattr(email_pipeline.send_email, "delay", m)
    return m

def _synthetic_wav(duration_s: float = 1.0, sample_rate: int = 48000) -> bytes:
    """Return WAV bytes — silent, 16-bit mono — that wave / ffprobe can read."""
    buf = io.BytesIO()
    n = int(sample_rate * duration_s)
    silence = np.zeros(n, dtype=np.int16)
    with wave.open(buf, "wb") as w:
        w.setnchannels(1)
        w.setsampwidth(2)
        w.setframerate(sample_rate)
        w.writeframes(silence.tobytes())
    return buf.getvalue()


@pytest.fixture()
def synthetic_wav_bytes() -> bytes:
    return _synthetic_wav(duration_s=1.0)


@pytest.fixture()
def mock_tts(monkeypatch: pytest.MonkeyPatch) -> dict[str, int]:
    """Patch httpx.get inside tts_service to return a synthetic WAV. Returns
    a counter dict the test can read to assert how many calls were made."""
    from app.services import tts_service as tts_mod

    counter = {"calls": 0}
    wav_bytes = _synthetic_wav(0.5)

    class _Resp:
        status_code = 200
        content = wav_bytes

        def raise_for_status(self) -> None:
            return None

    def _fake_get(url: str, params: dict[str, Any] | None = None, timeout: int = 0) -> _Resp:
        counter["calls"] += 1
        return _Resp()

    monkeypatch.setattr(tts_mod.httpx, "get", _fake_get)
    return counter


@pytest.fixture()
def mock_llm_split(monkeypatch: pytest.MonkeyPatch) -> dict[str, Any]:
    """Patch llm_service.split_and_annotate_ssml — return as many chunks as
    asked for, no warning. Tests override the return value when needed.
    """
    from app.services import llm_service as llm_mod

    state: dict[str, Any] = {"calls": 0, "chunks": None, "raise": None}

    async def _fake_split(
        script: str, slides_count: int, slide_texts: list[str] | None = None
    ) -> tuple[list[str], str | None]:
        state["calls"] += 1
        if state["raise"] is not None:
            raise state["raise"]
        chunks = state["chunks"]
        if chunks is None:
            chunks = [f"<p>chunk {i + 1}</p>" for i in range(slides_count)]
        return chunks, None

    monkeypatch.setattr(
        llm_mod.llm_service, "split_and_annotate_ssml", _fake_split
    )
    return state


@pytest.fixture()
def mock_vision(monkeypatch: pytest.MonkeyPatch) -> dict[str, Any]:
    """Patch vision_analysis_service.analyze_slide / analyze_presentation /
    summarize_presentation AND llm_service.refine_slide_narration so the
    pipelines complete without a real LLM call.
    """
    from app.services import llm_service as llm_mod
    from app.services import vision_analysis as vis_mod

    state: dict[str, Any] = {
        "analyze_calls": 0,
        "summarize_calls": 0,
        "analyze_raise": None,
        "analyze_return": "vision narration text",
    }

    async def _analyze_slide(
        slide_image_path: str,
        slide_number: int,
        total_slides: int,
        course_title: str,
        previous_context: str = "",
    ) -> str:
        state["analyze_calls"] += 1
        if state["analyze_raise"] is not None:
            raise state["analyze_raise"]
        return state["analyze_return"]

    async def _analyze_presentation(
        slide_image_paths: list[str],
        course_title: str,
        progress_cb: Any = None,
        cancel_check: Any = None,
    ) -> list[str]:
        # Mirror the real per-slide-boundary cancellation contract.
        if cancel_check is not None and cancel_check():
            raise vis_mod.AnalysisCancelled(0)
        if state["analyze_raise"] is not None:
            raise state["analyze_raise"]
        return [state["analyze_return"]] * len(slide_image_paths)

    async def _summarize_presentation(
        slide_image_paths: list[str], progress_cb: Any = None
    ) -> list[str]:
        state["summarize_calls"] += 1
        return [f"summary {i + 1}" for i in range(len(slide_image_paths))]

    async def _refine_slide_narration(vision_text: str, model: str | None = None) -> str:
        # Pass the vision text through unchanged — tests assert on the vision return value.
        return vision_text

    monkeypatch.setattr(vis_mod.vision_analysis_service, "analyze_slide", _analyze_slide)
    monkeypatch.setattr(
        vis_mod.vision_analysis_service, "analyze_presentation", _analyze_presentation
    )
    monkeypatch.setattr(
        vis_mod.vision_analysis_service, "summarize_presentation", _summarize_presentation
    )
    monkeypatch.setattr(llm_mod.llm_service, "refine_slide_narration", _refine_slide_narration)
    return state


@pytest.fixture()
def mock_subprocess(monkeypatch: pytest.MonkeyPatch) -> dict[str, list[list[str]]]:
    """Patch subprocess.run inside video_service so LibreOffice / pdftoppm /
    FFmpeg / ffprobe are not actually invoked. Creates a dummy PNG for each
    pdftoppm call and writes an empty file at any FFmpeg output path that
    looks like the last argument.
    """
    from app.services import video_service as vs_mod

    calls: dict[str, list[list[str]]] = {"all": []}

    class _Completed:
        def __init__(self, returncode: int = 0, stdout: str = "", stderr: bytes = b"") -> None:
            self.returncode = returncode
            self.stdout = stdout
            self.stderr = stderr

    def _fake_run(cmd: list[str], **kwargs: Any) -> _Completed:
        calls["all"].append(cmd)
        prog = Path(cmd[0]).name

        if prog == "libreoffice":
            # Args: --convert-to pdf --outdir <pdf_dir> <pptx>
            outdir = cmd[cmd.index("--outdir") + 1]
            src = cmd[-1]
            pdf_name = Path(src).stem + ".pdf"
            Path(outdir).mkdir(parents=True, exist_ok=True)
            (Path(outdir) / pdf_name).write_bytes(b"%PDF-1.4\n%%EOF\n")
            return _Completed()

        if prog == "pdftoppm":
            # Args: ... <pdf> <out_prefix> — out_prefix is last arg
            out_prefix = cmd[-1]
            out_dir = Path(out_prefix).parent
            out_dir.mkdir(parents=True, exist_ok=True)
            # Emit two PNGs so multi-slide pipelines have something to work with
            for i in (1, 2):
                (out_dir / f"slide-{i}.png").write_bytes(b"\x89PNG\r\n\x1a\n")
            return _Completed()

        if prog == "ffprobe":
            return _Completed(returncode=0, stdout="1.5\n")

        if prog == "ffmpeg":
            # Output path is always the last positional arg after the last -i.
            out = Path(cmd[-1])
            out.parent.mkdir(parents=True, exist_ok=True)
            out.write_bytes(_synthetic_wav(0.5))
            return _Completed()

        return _Completed()

    monkeypatch.setattr(vs_mod.subprocess, "run", _fake_run)
    return calls


# ── 7. Preserve pytest-asyncio's session loop across `asyncio.run` callers ─
# Two things call `asyncio.run(...)` mid-session, and its cleanup does
# `set_event_loop(None)`, wiping the thread's current loop:
#   • Celery pipelines inside eager-mode tasks
#     (video_pipeline._split_and_annotate) — i.e. during the test body;
#   • alembic/env.py (module level), triggered by the session-scoped
#     _alembic_upgraded fixture the first time a DB-backed test runs —
#     i.e. during fixture SETUP.
# pytest-asyncio runs every async test on the thread's CURRENT loop and only
# sets it when its scoped loop fixture is first created, so one wipe leaves
# every async test that runs after with "no current event loop in
# 'MainThread'". A sync snapshot/restore fixture can't fix this: it runs
# after session fixtures (missing the alembic wipe), and its snapshot may
# auto-create a stray loop that is NOT the loop async fixtures run on.
# Being async, this fixture executes ON pytest-asyncio's (session-scoped)
# loop itself, so it re-asserts that exact loop as the thread's current loop
# both before and after each test.

@pytest_asyncio.fixture(autouse=True)
async def _preserve_thread_event_loop() -> AsyncIterator[None]:
    import asyncio

    loop = asyncio.get_running_loop()
    asyncio.set_event_loop(loop)
    try:
        yield
    finally:
        if not loop.is_closed():
            asyncio.set_event_loop(loop)


# ── 8. Disable rate limit for unit-style fixtures (autouse safety net) ──────

@pytest.fixture(autouse=True)
def _disable_rate_limit() -> Iterator[None]:
    """Tests rarely care about throttling, but slowapi's limiter is shared
    module-global state. Disabling it here is cheap insurance against
    parametrized tests hitting 429."""
    try:
        from app.limiter import limiter

        prior = limiter.enabled
        limiter.enabled = False
        yield
        limiter.enabled = prior
    except Exception:
        yield


# ── 9. Sample PPTX fixture (session-scope, built once) ──────────────────────

@pytest.fixture(scope="session")
def sample_pptx_bytes() -> bytes:
    """Build a 2-slide PPTX in memory via python-pptx and return the raw bytes."""
    from pptx import Presentation

    prs = Presentation()
    blank = prs.slide_layouts[6]
    for i in range(2):
        slide = prs.slides.add_slide(blank)
        tx = slide.shapes.add_textbox(left=914400, top=914400, width=4572000, height=1828800)
        tx.text_frame.text = f"Slide {i + 1} content"

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


@pytest.fixture(scope="session")
def sample_pptx(fixtures_dir: Path, sample_pptx_bytes: bytes) -> Path:
    p = fixtures_dir / "sample.pptx"
    p.write_bytes(sample_pptx_bytes)
    return p


# ── Re-export for legacy use ─────────────────────────────────────────────────

__all__ = [
    "client",
    "app",
    "db_session",
    "teacher_user",
    "student_user",
    "teacher_token",
    "student_token",
    "mock_tts",
    "mock_llm_split",
    "mock_vision",
    "mock_subprocess",
    "sample_pdf",
    "sample_pptx",
    "sample_pptx_bytes",
    "synthetic_wav_bytes",
]
