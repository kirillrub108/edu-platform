"""Synthetic deck reproducing the LibreOffice no-wrap overflow.

The structure is copied from a real broken upload, brand-free and with neutral
text, so the regression test can live in the repo:

* three text boxes stacked flush, each exactly one line tall;
* ``<a:bodyPr wrap="none">`` plus ``<a:spAutoFit/>``;
* runs carry ``<a:latin typeface=""/>`` — the real typeface is declared only in
  the shape's ``lstStyle`` and names a font that resolves nowhere.

PowerPoint honours ``wrap="none"`` and never wraps. LibreOffice ignores it,
re-wraps the text in whatever font it substituted, declines to grow the
``spAutoFit`` box, and the overflowing line lands on the box below.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# Deliberately fictional: the point is a family neither embedded in the deck nor
# installed in the image, which is what forces LibreOffice to substitute.
MISSING_TYPEFACE = "AcmeSans-Regular"

# Sized like the real deck's: comfortably one line in the font the author used,
# just over the box once LibreOffice substitutes a wider one.
LINES = (
    "Строка один — 18-26 ноября",
    "Строка два — центральный зал",
    "Строка три — улица Ленина, 14",
)

_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
# One line of 20pt text plus the default 0.05" top/bottom insets.
_BOX_HEIGHT = Emu(400110)
_BOX_WIDTH = Inches(3.6)
_FIRST_TOP = Inches(1.0)


def _lst_style_xml(algn: str | None) -> str:
    """lstStyle declaring the unresolvable typeface, as PowerPoint writes it."""
    algn_attr = f' algn="{algn}"' if algn else ""
    return (
        f'<a:lstStyle xmlns:a="{_A}">'
        f"<a:lvl1pPr{algn_attr}>"
        f'<a:defRPr sz="2450" spc="50"><a:latin typeface="{MISSING_TYPEFACE}"/></a:defRPr>'
        f"</a:lvl1pPr></a:lstStyle>"
    )


def build_nowrap_deck(
    path: str | Path,
    *,
    algn: str | None = None,
    wrap: bool = False,
) -> Path:
    """Write the fixture deck to *path* and return it.

    ``algn`` goes on the shape's ``lvl1pPr`` (``"ctr"``/``"r"``/None for left),
    mirroring how the real deck declares alignment. ``wrap=True`` builds the
    control case — an ordinary wrapping box the pre-processing must leave alone.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    for index, text in enumerate(LINES):
        top = Emu(_FIRST_TOP + _BOX_HEIGHT * index)  # flush stack, no gaps
        box = slide.shapes.add_textbox(Inches(0.5), top, _BOX_WIDTH, _BOX_HEIGHT)
        text_frame = box.text_frame
        text_frame.word_wrap = wrap
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

        body = text_frame._txBody
        body.insert(1, parse_xml(_lst_style_xml(algn)))

        run = text_frame.paragraphs[0].add_run()
        run.text = text
        run.font.size = Pt(20)
        # Empty typeface: PowerPoint resolves it through lstStyle, LibreOffice
        # does not and falls back to its own default font instead.
        run.font._rPr.append(parse_xml(f'<a:latin xmlns:a="{_A}" typeface=""/>'))

    path = Path(path)
    prs.save(str(path))
    return path


def slide_boxes(path: str | Path) -> list[tuple[int, int]]:
    """``(left, width)`` in EMU for every text box on the first slide."""
    prs = Presentation(str(path))
    return [
        (int(shape.left), int(shape.width))
        for shape in prs.slides[0].shapes
        if shape.has_text_frame
    ]


def has_nowrap_body(path: str | Path) -> bool:
    """True when the deck still declares at least one ``wrap="none"`` box."""
    prs = Presentation(str(path))
    return any(
        shape.text_frame._txBody.find(qn("a:bodyPr")).get("wrap") == "none"
        for shape in prs.slides[0].shapes
        if shape.has_text_frame
    )
