"""Unit tests for the PPTX no-wrap pre-processing in app.services.video_service.

LibreOffice ignores ``<a:bodyPr wrap="none">`` and re-wraps text it rendered in a
substituted font, spilling the extra line onto the shape below. We widen those
boxes before conversion. These tests cover the geometry maths and the rewrite
itself; none of them need LibreOffice. The one test that does is marked ``slow``.
"""

from __future__ import annotations

import hashlib
import subprocess
from pathlib import Path
from typing import Any

import pytest
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches

from app.services import video_service as vs_mod
from app.services.video_service import (
    _pptx_cache_key,
    _prepare_pptx_for_libreoffice,
    _text_alignment,
    _widened_geometry,
)
from tests.fixtures.nowrap_deck import (
    MISSING_TYPEFACE,
    build_nowrap_deck,
    has_nowrap_body,
    slide_boxes,
)

pytestmark = pytest.mark.unit

SLIDE_CX = int(Inches(13.333))


class _RecordingLogger:
    """Stand-in for the structlog logger that just remembers the events."""

    def __init__(self) -> None:
        self.events: list[tuple[str, str, dict[str, Any]]] = []

    def _record(self, level: str):
        def log(event: str, **kw: Any) -> None:
            self.events.append((level, event, kw))

        return log

    def __getattr__(self, level: str):
        return self._record(level)

    def names(self) -> list[str]:
        return [event for _, event, _ in self.events]


@pytest.fixture
def log(monkeypatch: pytest.MonkeyPatch) -> _RecordingLogger:
    recorder = _RecordingLogger()
    monkeypatch.setattr(vs_mod, "logger", recorder)
    return recorder


# ── _widened_geometry ────────────────────────────────────────────────────────


def test_left_aligned_grows_rightwards_leaving_x_alone() -> None:
    left, width = int(Inches(1)), int(Inches(3))
    result = _widened_geometry(left, width, SLIDE_CX, PP_ALIGN.LEFT)

    assert result is not None
    new_left, new_width = result
    assert new_left == left  # ink starts where the author put it
    assert new_width > width


def test_right_aligned_grows_leftwards_pinning_the_right_edge() -> None:
    left, width = int(Inches(6)), int(Inches(3))
    result = _widened_geometry(left, width, SLIDE_CX, PP_ALIGN.RIGHT)

    assert result is not None
    new_left, new_width = result
    assert new_left < left
    assert new_left + new_width == left + width  # right edge pinned


def test_centred_grows_around_its_own_centre_not_the_slides() -> None:
    left, width = int(Inches(1)), int(Inches(3))
    centre = left + width / 2
    result = _widened_geometry(left, width, SLIDE_CX, PP_ALIGN.CENTER)

    assert result is not None
    new_left, new_width = result
    assert new_width > width
    assert new_left + new_width / 2 == pytest.approx(centre, abs=2)


def test_centred_near_the_edge_stays_centred_after_clamping() -> None:
    """Clamping must shrink both sides, never slide the centre sideways."""
    width = int(Inches(3))
    left = int(Inches(0.5))  # less room on the left than the factor wants
    centre = left + width / 2
    result = _widened_geometry(left, width, SLIDE_CX, PP_ALIGN.CENTER)

    assert result is not None
    new_left, new_width = result
    assert new_left >= 0
    assert new_left + new_width <= SLIDE_CX
    assert new_left + new_width / 2 == pytest.approx(centre, abs=2)


def test_left_aligned_box_already_at_the_slide_edge_cannot_grow() -> None:
    left = SLIDE_CX - int(Inches(3))
    assert _widened_geometry(left, int(Inches(3)), SLIDE_CX, PP_ALIGN.LEFT) is None


def test_right_aligned_box_already_at_the_left_edge_cannot_grow() -> None:
    assert _widened_geometry(0, int(Inches(3)), SLIDE_CX, PP_ALIGN.RIGHT) is None


def test_full_width_centred_box_cannot_grow() -> None:
    assert _widened_geometry(0, SLIDE_CX, SLIDE_CX, PP_ALIGN.CENTER) is None


# ── _text_alignment ──────────────────────────────────────────────────────────


def _first_text_frame(path: Path):
    from pptx import Presentation

    return Presentation(str(path)).slides[0].shapes[0].text_frame


@pytest.mark.parametrize(
    ("algn", "expected"),
    [(None, PP_ALIGN.LEFT), ("ctr", PP_ALIGN.CENTER), ("r", PP_ALIGN.RIGHT)],
)
def test_alignment_is_read_from_the_shape_lst_style(
    tmp_path: Path, algn: str | None, expected: PP_ALIGN
) -> None:
    """The real deck declares alignment only in lstStyle, never on the runs."""
    deck = build_nowrap_deck(tmp_path / "deck.pptx", algn=algn)
    assert _text_alignment(_first_text_frame(deck)) is expected


def test_explicit_paragraph_alignment_wins_over_lst_style(tmp_path: Path) -> None:
    deck = build_nowrap_deck(tmp_path / "deck.pptx", algn="ctr")
    text_frame = _first_text_frame(deck)
    text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    assert _text_alignment(text_frame) is PP_ALIGN.RIGHT


def test_mixed_paragraph_alignment_falls_back_to_left(tmp_path: Path) -> None:
    """Left is the only direction that provably cannot move the visible text."""
    deck = build_nowrap_deck(tmp_path / "deck.pptx")
    text_frame = _first_text_frame(deck)
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.add_paragraph().alignment = PP_ALIGN.RIGHT

    assert _text_alignment(text_frame) is PP_ALIGN.LEFT


# ── _prepare_pptx_for_libreoffice ────────────────────────────────────────────


def test_nowrap_boxes_are_widened_into_a_copy(tmp_path: Path, log: _RecordingLogger) -> None:
    deck = build_nowrap_deck(tmp_path / "deck.pptx")
    before = slide_boxes(deck)

    prepared = _prepare_pptx_for_libreoffice(str(deck), str(tmp_path / "work"))

    assert prepared != str(deck)
    after = slide_boxes(prepared)
    assert len(after) == len(before)
    for (_, old_width), (_, new_width) in zip(before, after, strict=True):
        assert new_width > old_width
    assert "pptx_nowrap_relaxed" in log.names()


def test_the_original_deck_is_left_untouched(tmp_path: Path, log: _RecordingLogger) -> None:
    """The cache key is hashed from the original bytes — see docs/DECISIONS.md §20."""
    deck = build_nowrap_deck(tmp_path / "deck.pptx")
    digest = hashlib.md5(deck.read_bytes()).hexdigest()
    key = _pptx_cache_key(str(deck))

    _prepare_pptx_for_libreoffice(str(deck), str(tmp_path / "work"))

    assert hashlib.md5(deck.read_bytes()).hexdigest() == digest
    assert _pptx_cache_key(str(deck)) == key


def test_the_prepared_copy_keeps_the_basename(tmp_path: Path, log: _RecordingLogger) -> None:
    """convert_pptx_to_images derives the produced PDF's name from this stem."""
    deck = build_nowrap_deck(tmp_path / "lecture one.pptx")

    prepared = _prepare_pptx_for_libreoffice(str(deck), str(tmp_path / "work"))

    assert Path(prepared).name == deck.name


def test_wrap_none_is_preserved_so_powerpoint_still_never_wraps(
    tmp_path: Path, log: _RecordingLogger
) -> None:
    deck = build_nowrap_deck(tmp_path / "deck.pptx")

    prepared = _prepare_pptx_for_libreoffice(str(deck), str(tmp_path / "work"))

    assert has_nowrap_body(prepared)


def test_an_ordinary_wrapping_deck_is_returned_unchanged(
    tmp_path: Path, log: _RecordingLogger
) -> None:
    """Regression guard: decks without the bug must render byte-identically."""
    deck = build_nowrap_deck(tmp_path / "deck.pptx", wrap=True)

    prepared = _prepare_pptx_for_libreoffice(str(deck), str(tmp_path / "work"))

    assert prepared == str(deck)
    assert "pptx_nowrap_relaxed" not in log.names()


def test_a_box_with_no_room_is_logged_and_skipped(tmp_path: Path, log: _RecordingLogger) -> None:
    """Cannot be fixed, must not be silent — the documented edge case."""
    from pptx import Presentation

    deck = build_nowrap_deck(tmp_path / "deck.pptx")
    prs = Presentation(str(deck))
    for shape in prs.slides[0].shapes:
        shape.left, shape.width = Emu(0), Emu(prs.slide_width)
    prs.save(str(deck))

    prepared = _prepare_pptx_for_libreoffice(str(deck), str(tmp_path / "work"))

    assert prepared == str(deck)
    assert "nowrap_widen_no_room" in log.names()


def test_a_corrupt_deck_degrades_to_the_original(tmp_path: Path, log: _RecordingLogger) -> None:
    """A layout tweak must never be the reason a lesson fails to render."""
    broken = tmp_path / "broken.pptx"
    broken.write_bytes(b"PK\x03\x04 definitely not a presentation")

    prepared = _prepare_pptx_for_libreoffice(str(broken), str(tmp_path / "work"))

    assert prepared == str(broken)
    assert "pptx_prepare_open_failed" in log.names()


# ── _log_missing_fonts ───────────────────────────────────────────────────────


def test_missing_typefaces_are_reported(
    tmp_path: Path, log: _RecordingLogger, monkeypatch: pytest.MonkeyPatch
) -> None:
    deck = build_nowrap_deck(tmp_path / "deck.pptx")
    monkeypatch.setattr(vs_mod, "_installed_font_families", lambda: frozenset({"calibri"}))

    vs_mod._log_missing_fonts(str(deck))

    reported = [kw for _, event, kw in log.events if event == "pptx_fonts_missing"]
    assert reported and MISSING_TYPEFACE in reported[0]["fonts"]


def test_theme_references_are_not_reported_as_missing(
    tmp_path: Path, log: _RecordingLogger, monkeypatch: pytest.MonkeyPatch
) -> None:
    """ "+mn-lt" & co. are resolved by the theme's own fontScheme, not fontconfig."""
    deck = build_nowrap_deck(tmp_path / "deck.pptx")
    monkeypatch.setattr(vs_mod, "_installed_font_families", lambda: frozenset({"calibri"}))

    vs_mod._log_missing_fonts(str(deck))

    reported = [kw for _, event, kw in log.events if event == "pptx_fonts_missing"]
    assert reported
    assert not [font for font in reported[0]["fonts"] if font.startswith("+")]


def test_font_scan_is_skipped_when_fontconfig_is_unavailable(
    tmp_path: Path, log: _RecordingLogger, monkeypatch: pytest.MonkeyPatch
) -> None:
    deck = build_nowrap_deck(tmp_path / "deck.pptx")
    monkeypatch.setattr(vs_mod, "_installed_font_families", frozenset)

    vs_mod._log_missing_fonts(str(deck))

    assert "pptx_fonts_missing" not in log.names()


# ── real LibreOffice ─────────────────────────────────────────────────────────


@pytest.mark.slow
def test_libreoffice_renders_the_fixture_without_overlapping_lines(tmp_path: Path) -> None:
    """End-to-end guard: the three lines must land on three separate rows.

    Without the pre-processing LibreOffice re-wraps each box and the overflow
    lands on the box below, which shows up here as fewer distinct text rows than
    there are boxes.
    """
    from lxml import etree

    def rows(pdf: Path) -> list[float]:
        xml = subprocess.run(
            ["pdftotext", "-bbox", "-f", "1", "-l", "1", str(pdf), "-"],
            capture_output=True,
            check=True,
        ).stdout
        root = etree.fromstring(xml)
        tops = {
            round(float(word.get("yMin")), 1)
            for word in root.iter("{http://www.w3.org/1999/xhtml}word")
        }
        return sorted(tops)

    def render(deck: str, out: Path) -> Path:
        out.mkdir(parents=True, exist_ok=True)
        subprocess.run(
            [
                "libreoffice",
                "--headless",
                f"-env:UserInstallation=file://{out}/profile",
                "--convert-to",
                "pdf",
                "--outdir",
                str(out),
                deck,
            ],
            capture_output=True,
            check=True,
            timeout=180,
        )
        return out / (Path(deck).stem + ".pdf")

    deck = build_nowrap_deck(tmp_path / "deck.pptx")
    baseline = rows(render(str(deck), tmp_path / "before"))

    prepared = _prepare_pptx_for_libreoffice(str(deck), str(tmp_path / "work"))
    fixed = rows(render(prepared, tmp_path / "after"))

    assert len(baseline) > 3, "fixture no longer reproduces the wrap — rewrite it"
    assert len(fixed) == 3, f"expected one row per box, got {fixed}"
